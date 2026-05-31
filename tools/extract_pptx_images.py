"""Extract images from a .pptx, mapping by slide number.

Usage:
    uv run python tools/extract_pptx_images.py PPTX OUT_DIR [--slides N N ...]

Output filenames are `pNN-iI.<ext>` where NN is the 1-based slide number and
I is the picture-shape index within that slide. If --slides is omitted, every
slide is processed.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract(pptx_path: Path, out_dir: Path, slide_numbers: list[int] | None) -> list[Path]:
    prs = Presentation(pptx_path)
    out_dir.mkdir(parents=True, exist_ok=True)
    targets = slide_numbers if slide_numbers else list(range(1, len(prs.slides) + 1))
    written: list[Path] = []
    for n in targets:
        if not 1 <= n <= len(prs.slides):
            print(f"  skip: slide {n} out of range (1..{len(prs.slides)})")
            continue
        slide = prs.slides[n - 1]
        for i, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                out_path = out_dir / f"p{n:02d}-i{i}.{img.ext}"
                out_path.write_bytes(img.blob)
                written.append(out_path)
                print(f"  saved: {out_path.name} ({len(img.blob)} bytes)")
    return written


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="source .pptx")
    parser.add_argument("out", type=Path, help="output directory")
    parser.add_argument(
        "--slides", type=int, nargs="*", default=None,
        help="1-based slide numbers to extract; omit for all",
    )
    args = parser.parse_args()
    written = extract(args.pptx, args.out, args.slides)
    print(f"done: {len(written)} image(s) written to {args.out}")


if __name__ == "__main__":
    main()
